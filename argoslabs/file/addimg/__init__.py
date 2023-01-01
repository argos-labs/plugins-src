#!/usr/bin/env python
# coding=utf8
"""
====================================
 :mod:`argoslabs.file.addimg`
====================================
.. moduleauthor:: Irene Cho <irene@argos-labs.com>
.. note:: ARGOS-LABS License

Description
===========
ARGOS LABS plugin module for adding an image to the file
"""
#
# Authors
# ===========
#
# * Irene Cho
#
# Change Log
# --------
#
#  * [2020/10/02]
#     - build a plugin
#  * [2020/10/02]
#     - starting
#

################################################################################
import os
import csv
import sys
import shutil
import openpyxl
from docx import Document
from pptx import Presentation
# from docx.shared import Inches, Cm
from alabs.common.util.vvencoding import get_file_encoding
from reportlab.pdfgen import canvas
from PyPDF2 import PdfFileWriter, PdfFileReader
from alabs.common.util.vvargs import func_log, get_icon_path, ModuleContext, \
    ArgsError, ArgsExit


################################################################################
# noinspection PyBroadException
class addimgAPI(object):

    # ==========================================================================
    def __init__(self, argspec):
        self.argspec = argspec
        self.filename = argspec.filename
        self.img = argspec.img
        for i in [self.filename, self.img]:
            if not os.path.exists(i):
                raise IOError('Cannot read the file "%s"' % i)
        self.output = argspec.output
        if not self.output:
            self.output = self._get_safe_next_filename(self.filename)
        self.dx = argspec.dx
        self.dy = argspec.dy
        self.dwidth = argspec.dwidth
        self.dheight = argspec.dheight
        self.tempfile = None

    # ==========================================================================
    @staticmethod
    def _get_safe_next_filename(fn):
        fn, ext = os.path.splitext(fn)
        for n in range(1, 1000000):
            nfn = f'{fn} ({n})' + ext
            if not os.path.exists(nfn):
                return nfn

    # ==========================================================================
    def addimgpdf(self):
        self.tempfile = self._get_safe_next_filename('temp.pdf')
        c = canvas.Canvas(self.tempfile)
        if self.dwidth != 0 and self.dheight != 0:
            c.drawImage(self.img, self.dx, self.dy, self.dwidth,
                        self.dheight,
                        mask=(0, 0, 0, 0, 0, 0))  # position, scale
        else:
            c.drawImage(self.img, self.dx, self.dy)
        c.save()
        # The mask parameter lets you create transparent images. It takes 6 numbers and defines the range of RGB values which will be masked out or treated as transparent. For example with [0,2,40,42,136,139], it will mask
        # out any pixels with a Red value from 0 or 1, Green from 40 or 41 and Blue of 136, 137 or 138 (on a scale of
        # 0-255). It's currently your job to know which color is the 'transparent' or background one.
        t = open(self.tempfile, "rb")
        watermark = PdfFileReader(t)
        output_file = PdfFileWriter()
        input_file = PdfFileReader(open(self.filename, "rb"))
        page_count = input_file.getNumPages()
        pn = self.argspec.pagenum - 1
        if page_count - 1 < pn:
            raise IOError('The page number exceeds the number of pages')
        for i in range(page_count):
            input_page = input_file.getPage(i)
            if i == pn:
                input_page.mergePage(watermark.getPage(0))
            output_file.addPage(input_page)
        with open(self.output, "wb") as outputStream:
            output_file.write(outputStream)
        if self.tempfile and os.path.exists(self.tempfile):
            t.close()
            os.remove(self.tempfile)

    # ==========================================================================
    def addimgdocx(self):
        document = Document(self.filename)
        if self.dwidth and self.dheight and self.dwidth != 0 and self.dheight != 0:
            document.add_picture(self.img, self.dwidth * 10000,
                                 self.dheight * 10000)
        elif self.dwidth and self.dwidth != 0:
            document.add_picture(self.img, self.dwidth * 10000)
        elif self.dheight and self.dheight != 0:
            document.add_picture(self.img, self.dheight * 10000)
        else:
            document.add_picture(self.img)
        # scale
        document.save(self.output)

    # ==========================================================================
    def addimgppt(self):
        prs = Presentation(self.filename)
        slide = prs.slides[0]
        if self.dwidth and self.dheight and self.dwidth != 0 and self.dheight != 0:
            slide.shapes.add_picture(self.img, self.dx * 10000, self.dx * 10000,
                                     width=self.dwidth * 10000,
                                     height=self.dheight * 10000)
        elif self.dwidth and self.dwidth != 0:
            slide.shapes.add_picture(self.img, self.dx * 10000, self.dx * 10000,
                                     width=self.dwidth * 10000)
        elif self.dheight and self.dheight != 0:
            slide.shapes.add_picture(self.img, self.dx * 10000, self.dx * 10000,
                                     height=self.dheight * 10000)
        else:
            slide.shapes.add_picture(self.img, self.dx * 10000,
                                     self.dx * 10000, )
        prs.save(self.output)

    # ==========================================================================
    def addimgexcel(self):
        if self.argspec.dataonly:
            data_only = True
            shutil.copy(self.filename, self.output)
            self.filename = self.output
        else:
            data_only = False
        wb = openpyxl.load_workbook(self.filename, read_only=False,
                                    data_only=data_only, keep_vba=False)
        if self.argspec.sheetname:
            try:
                ws = wb[self.argspec.sheetname]
            except Exception:
                wb.create_sheet(self.argspec.sheetname)
                ws = wb[self.argspec.sheetname]
        else:
            ws = wb.active
        p = open(self.img, 'rb')
        img = openpyxl.drawing.image.Image(p)
        if self.dwidth and self.dwidth != 0:
            img.width = self.dwidth
        if self.dheight and self.dheight != 0:
            img.height = self.dheight
        ws.add_image(img, self.argspec.cell)
        wb.save(self.output)
        p.close()

    # ==========================================================================
    @staticmethod
    def csv2xls(s, t, encoding=None):
        wb = openpyxl.Workbook()
        ws = wb.active
        if not encoding:
            encoding = get_file_encoding(s)
        try:
            with open(s, encoding=encoding) as f:
                reader = csv.reader(f)
                for row in reader:
                    ws.append(row)
            wb.save(t)
        except Exception as e:
            raise RuntimeError('Cannot read CSV file. Please check encoding: %s'
                               % str(e))
        return True

    # ==========================================================================
    def do(self, fn):
        filetype = os.path.splitext(fn.lower())[-1]
        if filetype == '.pdf':
            self.addimgpdf()
        elif filetype == '.docx':
            self.addimgdocx()
        elif filetype == '.pptx':
            self.addimgppt()
        elif filetype in '.xlsx':
            self.addimgexcel()
        print(os.path.abspath(self.output), end='')


################################################################################
@func_log
def do_addimg(mcxt, argspec):
    mcxt.logger.info('>>>starting...')
    try:
        ai = addimgAPI(argspec)
        ai.do(argspec.filename)
        return 0
    except Exception as err:
        msg = str(err)
        mcxt.logger.error(msg)
        sys.stderr.write('%s%s' % (msg, os.linesep))
        # print(False)
        return 1
    finally:
        sys.stdout.flush()
        mcxt.logger.info('>>>end...')


################################################################################
def _main(*args):
    with ModuleContext(
            owner='ARGOS-LABS',
            group='9',
            version='1.0',
            platform=['windows', 'darwin', 'linux'],
            output_type='csv',
            display_name='Attach Image',
            icon_path=get_icon_path(__file__),
            description='Add an image to different documents',
    ) as mcxt:
        # ######################################### for app dependent options
        mcxt.add_argument('--output', display_name='Output File',
                          help='outputfile')
        # ----------------------------------------------------------------------
        mcxt.add_argument('--dx', display_name='Offset Horizontally',
                          type=float,
                          help='move an image to the left', default=0)
        # ----------------------------------------------------------------------
        mcxt.add_argument('--dy', display_name='Offset Vertically',
                          type=float,
                          help='move an image to upside in pdf or downside in ppt',
                          default=0)
        # ----------------------------------------------------------------------
        mcxt.add_argument('--dwidth', display_name='Resize Image(w)', type=int,
                          help='resize the image width')
        # ----------------------------------------------------------------------
        mcxt.add_argument('--dheight', display_name='Resize Image(h)', type=int,
                          help='resize the image height')
        # ----------------------------------------------------------------------
        mcxt.add_argument('--pagenum', display_name='PDF Page Number ',
                          type=int, default=1,
                          help='page number in pdf')
        # ----------------------------------------------------------------------
        mcxt.add_argument('--cell', display_name='Excel Paste Cell',
                          help='specify the cell to paste the image in Excel')
        # ----------------------------------------------------------------------
        mcxt.add_argument('--sheetname', display_name='Excel Paste Sheet',
                          help='sheet name to handle. If not specified last'
                               ' activated sheet is chosen')
        # ----------------------------------------------------------------------
        mcxt.add_argument('--dataonly', display_name='Excel Data Only',
                          action='store_true',
                          help='print only data not formulas')
        # ----------------------------------------------------------------------
        mcxt.add_argument('--encoding',
                          default='utf-8',
                          display_name='Excel Encoding',
                          help='Encoding for CSV file, default is [[utf-8]]')
        # ###################################### for app dependent parameters
        mcxt.add_argument('filename',
                          display_name='Filename', input_method='fileread',
                          help='PDF,docx,ppt,xlsx')
        mcxt.add_argument('img', display_name='Image',
                          input_method='fileread', help='jpg,png,jpeg')

        argspec = mcxt.parse_args(args)
        return do_addimg(mcxt, argspec)


################################################################################
def main(*args):
    try:
        return _main(*args)
    except ArgsError as err:
        sys.stderr.write('Error: %s\nPlease -h to print help\n' % str(err))
    except ArgsExit as _:
        pass
