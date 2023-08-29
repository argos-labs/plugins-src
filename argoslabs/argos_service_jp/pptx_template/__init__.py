#!/usr/bin/env python
# coding=utf8
"""
====================================
 :mod:`argoslabs.argos_service_jp.pptx_template`
====================================
.. moduleauthor:: Taiki Shimasaki <shimasaki@argos-service.com>
.. note:: ARGOS-LABS License

Description
===========
Input Plugin Description
"""
# Authors
# ===========
#
# * Taiki Shimasaki
#
# Change Log
# --------
#
#  * [2021/07/14]
#     Create

################################################################################
import os
import sys
from alabs.common.util.vvargs import ModuleContext, func_log, \
    ArgsError, ArgsExit, get_icon_path
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Cm, Inches, Pt
from PIL import Image
from pathvalidate import sanitize_filename
import re

################################################################################
class Operate_Template(object):

    # ==========================================================================
    def __init__(self, pp_path, file_exists):
        self.pp_path = os.path.abspath(pp_path)
        self.out_path = None

        self.prs = Presentation(pp_path)

        self.file_exists = file_exists

    # ==========================================================================
    def check_file_exists(self):
        if not os.path.isfile(self.pp_path):
            raise IOError('File {} not found'.format(self.pp_path))

        else:
            pass

    # ==========================================================================
    def file_type_check(self):
        temp_file_name = os.path.basename(self.pp_path)
        temp_name_split = temp_file_name.rsplit('.', 1)
        if len(temp_name_split) == 1:
            raise IOError('Cannot deal with the file!')
        else:
            temp_ext = temp_name_split[1]
            temp_ext = str.lower(temp_ext)

        if temp_ext != 'pptx':
            raise IOError('This file does not appear to be a PPTX file!')
        else:
            pass

    # ==========================================================================
    def img_path_check(self, img_path, img_exists):
        img_path = os.path.abspath(img_path)
        if not os.path.isfile(img_path):
            if img_exists == 'Return Failure':
                raise IOError('File {} not found'.format(img_path))
            elif img_exists == 'Ignore (Not Replace)':
                return 'ignore'

        else:
            pass

    # ==========================================================================
    def out_path_check(self, out_path):
        out_path = os.path.abspath(out_path)

        out_path_name = os.path.basename(out_path)
        out_path_name = sanitize_filename(out_path_name)
        out_path_dir = os.path.dirname(out_path)

        if not os.path.isdir(out_path_dir):
            raise IOError('Directory {} not found'.format(out_path_dir))

        else:
            pass

        out_name_split = out_path_name.rsplit('.', 1)
        out_file_name = out_name_split[0]
        if len(out_name_split) == 1:
            raise IOError('Output file extension is None!')
        else:
            out_ext = out_name_split[1]
            out_ext = str.lower(out_ext)

        if out_ext != 'pptx':
            raise IOError('Output file does not appear to be a PPTX file!')
        else:
            pass

        self.out_path = '{}{}{}.{}'.format(out_path_dir,
                                           os.sep,
                                           out_file_name,
                                           out_ext)

        # ----------------------------------------------------------------------
        if os.path.exists(self.out_path):
            if self.file_exists == 'Return Failure':
                raise IOError('File is already exists!')

            elif self.file_exists == 'Overwrite':
                pass
                # os.remove(self.out_path)

            elif self.file_exists == 'Add (n) at End':
                count = 0
                while os.path.exists(self.out_path):
                    out_file_name = re.sub(r'\([0-9]+\)$', '', out_file_name)
                    count += 1
                    out_file_name += '({})'.format(count)
                    self.out_path = '{}{}{}.{}'.format(out_path_dir,
                                                       os.sep,
                                                       out_file_name,
                                                       out_ext)
                else:
                    pass

            else:
                raise IOError('Unexpected Error Occurred')

        else:
            pass

    # ==========================================================================
    def replace_text(self, i, text):
        text_paragraph = self.prs.slides[0].shapes[i].text_frame.paragraphs[0]

        # ======================================================================
        def replace_paragraph(paragraph, new_text):
            p = paragraph._p
            for idx, run in enumerate(paragraph.runs):
                if idx == 0:
                    continue
                p.remove(run._r)
            paragraph.runs[0].text = new_text

        replace_paragraph(text_paragraph, text)

    # ==========================================================================
    def replace_img(self, i, img_path):
        img_path = os.path.abspath(img_path)
        slide = self.prs.slides[0]
        image = self.prs.slides[0].shapes[i]

        # ======================================================================
        def replace_picture(slide_obj, shape_obj, img):
            with open(img, "rb") as file_obj:
                r_img_blob = file_obj.read()
            img_pic = shape_obj._pic
            img_rid = img_pic.xpath("./p:blipFill/a:blip/@r:embed")[0]
            img_part = slide_obj.part.related_parts[img_rid]
            img_part._blob = r_img_blob

        replace_picture(slide, image, img_path)

    # ==========================================================================
    def save_file(self):
        self.prs.save(self.out_path)

        print(self.out_path, end="")

################################################################################
@func_log
def replace_object(mcxt, argspec):
    mcxt.logger.info('>>>starting...')
    try:
        OP = Operate_Template(argspec.pp_path, argspec.file_exists)
        OP.check_file_exists()
        OP.file_type_check()
        OP.out_path_check(argspec.out_path)

        # Select Template ------------------------------------------------------
        if argspec.template == 'x10':
            # Text =============================================================

            if argspec.text_01:
                OP.replace_text(19, argspec.text_01)
            else:
                pass

            if argspec.text_02:
                OP.replace_text(18, argspec.text_02)
            else:
                pass

            if argspec.text_03:
                OP.replace_text(17, argspec.text_03)
            else:
                pass

            if argspec.text_04:
                OP.replace_text(16, argspec.text_04)
            else:
                pass

            if argspec.text_05:
                OP.replace_text(15, argspec.text_05)
            else:
                pass

            if argspec.text_06:
                OP.replace_text(14, argspec.text_06)
            else:
                pass

            if argspec.text_07:
                OP.replace_text(13, argspec.text_07)
            else:
                pass

            if argspec.text_08:
                OP.replace_text(12, argspec.text_08)
            else:
                pass

            if argspec.text_09:
                OP.replace_text(11, argspec.text_09)
            else:
                pass

            if argspec.text_10:
                OP.replace_text(10, argspec.text_10)
            else:
                pass

            # Image ============================================================

            if argspec.img_01:
                if OP.img_path_check(argspec.img_01, argspec.img_exists) == 'ignore':
                    pass
                else:
                    OP.replace_img(9, argspec.img_01)
            else:
                pass

            if argspec.img_02:
                if OP.img_path_check(argspec.img_02, argspec.img_exists) == 'ignore':
                    pass
                else:
                    OP.replace_img(8, argspec.img_02)
            else:
                pass

            if argspec.img_03:
                if OP.img_path_check(argspec.img_03, argspec.img_exists) == 'ignore':
                    pass
                else:
                    OP.replace_img(7, argspec.img_03)
            else:
                pass

            if argspec.img_04:
                if OP.img_path_check(argspec.img_04, argspec.img_exists) == 'ignore':
                    pass
                else:
                    OP.replace_img(6, argspec.img_04)
            else:
                pass

            if argspec.img_05:
                if OP.img_path_check(argspec.img_05, argspec.img_exists) == 'ignore':
                    pass
                else:
                    OP.replace_img(5, argspec.img_05)
            else:
                pass

            if argspec.img_06:
                if OP.img_path_check(argspec.img_06, argspec.img_exists) == 'ignore':
                    pass
                else:
                    OP.replace_img(4, argspec.img_06)
            else:
                pass

            if argspec.img_07:
                if OP.img_path_check(argspec.img_07, argspec.img_exists) == 'ignore':
                    pass
                else:
                    OP.replace_img(3, argspec.img_07)
            else:
                pass

            if argspec.img_08:
                if OP.img_path_check(argspec.img_08, argspec.img_exists) == 'ignore':
                    pass
                else:
                    OP.replace_img(2, argspec.img_08)
            else:
                pass

            if argspec.img_09:
                if OP.img_path_check(argspec.img_09, argspec.img_exists) == 'ignore':
                    pass
                else:
                    OP.replace_img(1, argspec.img_09)
            else:
                pass

            if argspec.img_10:
                if OP.img_path_check(argspec.img_10, argspec.img_exists) == 'ignore':
                    pass
                else:
                    OP.replace_img(0, argspec.img_10)
            else:
                pass

        elif argspec.template == 'x6':
            # Text =============================================================

            if argspec.text_01:
                OP.replace_text(11, argspec.text_01)
            else:
                pass

            if argspec.text_02:
                OP.replace_text(10, argspec.text_02)
            else:
                pass

            if argspec.text_03:
                OP.replace_text(9, argspec.text_03)
            else:
                pass

            if argspec.text_04:
                OP.replace_text(8, argspec.text_04)
            else:
                pass

            if argspec.text_05:
                OP.replace_text(7, argspec.text_05)
            else:
                pass

            if argspec.text_06:
                OP.replace_text(6, argspec.text_06)
            else:
                pass

            # Image ============================================================

            if argspec.img_01:
                if OP.img_path_check(argspec.img_01,
                                     argspec.img_exists) == 'ignore':
                    pass
                else:
                    OP.replace_img(5, argspec.img_01)
            else:
                pass

            if argspec.img_02:
                if OP.img_path_check(argspec.img_02,
                                     argspec.img_exists) == 'ignore':
                    pass
                else:
                    OP.replace_img(4, argspec.img_02)
            else:
                pass

            if argspec.img_03:
                if OP.img_path_check(argspec.img_03,
                                     argspec.img_exists) == 'ignore':
                    pass
                else:
                    OP.replace_img(3, argspec.img_03)
            else:
                pass

            if argspec.img_04:
                if OP.img_path_check(argspec.img_04,
                                     argspec.img_exists) == 'ignore':
                    pass
                else:
                    OP.replace_img(2, argspec.img_04)
            else:
                pass

            if argspec.img_05:
                if OP.img_path_check(argspec.img_05,
                                     argspec.img_exists) == 'ignore':
                    pass
                else:
                    OP.replace_img(1, argspec.img_05)
            else:
                pass

            if argspec.img_06:
                if OP.img_path_check(argspec.img_06,
                                     argspec.img_exists) == 'ignore':
                    pass
                else:
                    OP.replace_img(0, argspec.img_06)
            else:
                pass

        elif argspec.template == 'x3':
            # Text =============================================================

            if argspec.text_01:
                OP.replace_text(5, argspec.text_01)
            else:
                pass

            if argspec.text_02:
                OP.replace_text(4, argspec.text_02)
            else:
                pass

            if argspec.text_03:
                OP.replace_text(3, argspec.text_03)
            else:
                pass

            # Image ============================================================

            if argspec.img_01:
                if OP.img_path_check(argspec.img_01,
                                     argspec.img_exists) == 'ignore':
                    pass
                else:
                    OP.replace_img(2, argspec.img_01)
            else:
                pass

            if argspec.img_02:
                if OP.img_path_check(argspec.img_02,
                                     argspec.img_exists) == 'ignore':
                    pass
                else:
                    OP.replace_img(1, argspec.img_02)
            else:
                pass

            if argspec.img_03:
                if OP.img_path_check(argspec.img_03,
                                     argspec.img_exists) == 'ignore':
                    pass
                else:
                    OP.replace_img(0, argspec.img_03)
            else:
                pass
        # ----------------------------------------------------------------------

        OP.save_file()

        mcxt.logger.info('>>>end...')
        return 0
    except Exception as err:
        msg = str(err)
        mcxt.logger.error(msg)
        sys.stderr.write('%s%s' % (msg, os.linesep))
        return 1
    finally:
        sys.stdout.flush()
        mcxt.logger.info('>>>end...')


################################################################################
def _main(*args):
    with ModuleContext(
        owner='ARGOS-SERVICE-JP',
        group='2',
        version='1.0',
        platform=['windows', 'darwin', 'linux'],
        output_type='text',
        display_name='PPTX Template',
        icon_path=get_icon_path(__file__),
        description='Replace Text & Image in Template PPTX File',
    ) as mcxt:
        # ##################################### for app dependent parameters
        mcxt.add_argument('template',
                          display_name='Template Type',
                          default='x10',
                          choices=['x10', 'x6', 'x3'],
                          help='Select the type of template based on the number of shapes')
        mcxt.add_argument('pp_path',
                          display_name='Template Path',
                          input_method='fileread',
                          help='Select, and input the downloaded \"Template-10.pptx\" path')
        mcxt.add_argument('out_path',
                          display_name='Output File Path',
                          input_method='fileread',
                          help='Input a output .pptx file path')
        # ######################################## for app dependent options
        mcxt.add_argument('--text_01',
                          display_name='Text01',
                          show_default=True,
                          input_group='Text',
                          help='Input the text to be replaced in Text01')
        mcxt.add_argument('--text_02',
                          display_name='Text02',
                          show_default=True,
                          input_group='Text',
                          help='Input the text to be replaced in Text02')
        mcxt.add_argument('--text_03',
                          display_name='Text03',
                          show_default=True,
                          input_group='Text',
                          help='Input the text to be replaced in Text03')
        mcxt.add_argument('--text_04',
                          display_name='Text04',
                          show_default=True,
                          input_group='Text',
                          help='Input the text to be replaced in Text04')
        mcxt.add_argument('--text_05',
                          display_name='Text05',
                          show_default=True,
                          input_group='Text',
                          help='Input the text to be replaced in Text05')
        mcxt.add_argument('--text_06',
                          display_name='Text06',
                          show_default=True,
                          input_group='Text',
                          help='Input the text to be replaced in Text06')
        mcxt.add_argument('--text_07',
                          display_name='Text07',
                          show_default=True,
                          input_group='Text',
                          help='Input the text to be replaced in Text07')
        mcxt.add_argument('--text_08',
                          display_name='Text08',
                          show_default=True,
                          input_group='Text',
                          help='Input the text to be replaced in Text08')
        mcxt.add_argument('--text_09',
                          display_name='Text09',
                          show_default=True,
                          input_group='Text',
                          help='Input the text to be replaced in Text09')
        mcxt.add_argument('--text_10',
                          display_name='Text10',
                          show_default=True,
                          input_group='Text',
                          help='Input the text to be replaced in Text10')
        mcxt.add_argument('--img_01',
                          display_name='Image01',
                          input_method="fileread",
                          show_default=True,
                          input_group='Image',
                          help='Select the image to be replaced in Image01')
        mcxt.add_argument('--img_02',
                          display_name='Image02',
                          input_method="fileread",
                          show_default=True,
                          input_group='Image',
                          help='Select the image to be replaced in Image02')
        mcxt.add_argument('--img_03',
                          display_name='Image03',
                          input_method="fileread",
                          show_default=True,
                          input_group='Image',
                          help='Select the image to be replaced in Image03')
        mcxt.add_argument('--img_04',
                          display_name='Image04',
                          input_method="fileread",
                          show_default=True,
                          input_group='Image',
                          help='Select the image to be replaced in Image04')
        mcxt.add_argument('--img_05',
                          display_name='Image05',
                          input_method="fileread",
                          show_default=True,
                          input_group='Image',
                          help='Select the image to be replaced in Image05')
        mcxt.add_argument('--img_06',
                          display_name='Image06',
                          input_method="fileread",
                          show_default=True,
                          input_group='Image',
                          help='Select the image to be replaced in Image06')
        mcxt.add_argument('--img_07',
                          display_name='Image07',
                          input_method="fileread",
                          show_default=True,
                          input_group='Image',
                          help='Select the image to be replaced in Image07')
        mcxt.add_argument('--img_08',
                          display_name='Image08',
                          input_method="fileread",
                          show_default=True,
                          input_group='Image',
                          help='Select the image to be replaced in Image08')
        mcxt.add_argument('--img_09',
                          display_name='Image09',
                          input_method="fileread",
                          show_default=True,
                          input_group='Image',
                          help='Select the image to be replaced in Image09')
        mcxt.add_argument('--img_10',
                          display_name='Image10',
                          input_method="fileread",
                          show_default=True,
                          input_group='Image',
                          help='Select the image to be replaced in Image10')
        mcxt.add_argument('--file_exists',
                          display_name='If File Exists',
                          default='Add (n) at End',
                          choices=['Add (n) at End',
                                   'Return Failure',
                                   'Overwrite'],
                          help='When output file already exists')
        mcxt.add_argument('--img_exists',
                          display_name='Image Not Exists',
                          default='Return Failure',
                          choices=['Return Failure',
                                   'Ignore (Not Replace)'],
                          help='When output file already exists')

        argspec = mcxt.parse_args(args)
        return replace_object(mcxt, argspec)


################################################################################
def main(*args):
    try:
        return _main(*args)
    except ArgsError as err:
        sys.stderr.write('Error: %s\nPlease -h to print help\n' % str(err))
    except ArgsExit as _:
        pass
